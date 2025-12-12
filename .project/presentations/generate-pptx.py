"""Generate PowerPoint presentation for Aura Developer Talk."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BG = RGBColor(30, 30, 30)
ACCENT = RGBColor(0, 120, 212)  # Azure blue
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullets, code=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Bullets
    if code:
        bullet_width = Inches(6)
    else:
        bullet_width = Inches(12.333)
    
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), bullet_width, Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)
    
    # Code block
    if code:
        code_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = LIGHT_GRAY
    
    return slide

def add_table_slide(title, headers, rows):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Table
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
    
    return slide

# ============ SLIDES ============

# Title slide
add_title_slide("Aura", "A Local-First AI Foundation for Developer Workflows")

# Part 1 title
add_title_slide("Part 1", "What is Aura?")

# The Hook
add_content_slide("The Hook", [
    "What if you could have 5 AI agents working on your codebase simultaneously?",
    "Each in its own git branch",
    "While you continue coding on main?",
    "",
    "That's Aura."
])

# The Problem
add_content_slide("The Problem", [
    "Context Switching — Every AI interaction breaks your flow",
    "The Copilot Ceiling — Single-turn completions can't handle multi-step tasks",
    "Privacy vs Quality — Local models are safe but weak, cloud is powerful but risky",
    "",
    "\"AI is fast at writing wrong code\""
])

# The Solution
add_content_slide("The Solution: Aura", [
    "Index locally, generate with the best model",
    "Your code stays on your machine",
    "Only queries go to cloud",
    "Swap providers with config change",
    "",
    "Architecture:",
    "  Foundation: Agents, LLM abstraction, RAG, Tools",
    "  Modules: Developer (now), Research (future), Personal (future)"
])

# Live Demo 1
add_content_slide("Demo: RAG Search", [
    "2,691 semantic chunks indexed",
    "30+ languages via TreeSitter",
    "Query: \"workflow execution\"",
    "",
    "Show: VS Code Aura sidebar"
], code="""curl http://localhost:5300/api/rag/search?query=workflow+execution

# Returns semantic matches from:
# - Python functions
# - TypeScript interfaces  
# - Rust structs
# - C# classes
# - And 26 more languages""")

# Live Demo 2
add_content_slide("Demo: Concurrent Workflows", [
    "Each workflow gets isolated git worktree",
    "No merge conflicts between workflows",
    "UI shows real-time progress",
    "",
    "Create → Plan → Execute → Review"
], code="""POST /api/developer/workflows
{
  "repositoryPath": "C:/work/aura",
  "issue": {
    "title": "Add health check",
    "body": "Add /health endpoint"
  }
}

# Creates:
# - Workflow record
# - Git worktree
# - Execution plan""")

# Part 2 title
add_title_slide("Part 2", "How We Built Aura")

# The Meta Challenge
add_content_slide("The Meta Challenge", [
    "We're building an AI system... using AI",
    "How do you avoid infinite recursion?",
    "",
    "The trap:",
    "  AI writes code fast →",
    "  Wrong code accumulates →",
    "  Rewrite everything"
])

# What Didn't Work
add_content_slide("What Didn't Work", [
    "\"Just let AI code\"",
    "  → Got 38,000 lines of inconsistent code",
    "",
    "\"Review everything\"",
    "  → Couldn't keep up with AI output",
    "",
    "Both approaches led to rewrites"
])

# What Worked
add_content_slide("What Worked: Spec-Driven Development", [
    "1. Human writes spec (requirements, constraints)",
    "2. AI reads spec, proposes implementation",
    "3. Human approves or adjusts",
    "4. AI implements",
    "5. Human verifies",
    "",
    "AI doesn't decide architecture — only implementation"
])

# Extensibility intro
add_content_slide("Extensibility Through Conversation", [
    "\"A lot of the flexibility came from conversation",
    "between the human and the Agent\"",
    "",
    "5 examples of human-AI negotiation →"
])

# Example 1
add_content_slide("1. Hot-Reload Agents", [
    "Human: \"I want to edit agents without restarting\"",
    "AI: \"Markdown files with FileSystemWatcher\"",
    "",
    "Result: Drop a .md file → agent is live"
], code="""# agents/coding-agent.md

# Coding Agent

## Metadata
- **Priority**: 70

## Capabilities
- coding
- testing

## System Prompt
You are an expert developer...""")

# Example 2
add_content_slide("2. Pluggable LLM Providers", [
    "Human: \"Switch between Ollama and OpenAI easily\"",
    "AI: \"ILlmProviderRegistry that resolves by config\"",
    "",
    "Result: One config change swaps providers"
], code=""""Llm": {
  "DefaultProvider": "AzureOpenAI",
  "Providers": {
    "Ollama": { 
      "DefaultModel": "qwen2.5-coder:7b" 
    },
    "OpenAI": { 
      "DefaultModel": "gpt-4o" 
    },
    "AzureOpenAI": { 
      "DefaultDeployment": "gpt-4o" 
    }
  }
}""")

# Example 3
add_content_slide("3. Externalized Prompts", [
    "Human: \"Prompts in code = hard to iterate\"",
    "AI: \"Handlebars templates? Hot-reload?\"",
    "Human: \"Yes, with front-matter for metadata\"",
    "",
    "Result: Edit prompt → see change immediately"
], code="""# prompts/workflow-plan.prompt
---
description: Creates execution plan
---
Create a plan for this task.

## Issue Title
{{title}}

## Issue Description  
{{description}}""")

# Example 4
add_content_slide("4. Module System", [
    "Human: \"Developer now, Research later, Personal eventually\"",
    "AI: \"IAuraModule interface — each registers its own services\"",
    "Human: \"But no cross-module dependencies!\"",
    "",
    "Result: Plug-in verticals with clean boundaries"
], code="""public class DeveloperModule : IAuraModule
{
    public string ModuleId => "developer";
    public string Name => "Developer Workflow";
    
    // Only depends on Foundation
    public IReadOnlyList<string> Dependencies => [];
    
    public void ConfigureServices(...)
    {
        services.AddScoped<IWorkflowService>();
        // ...
    }
}""")

# Example 5
add_content_slide("5. TreeSitter Multi-Language", [
    "Human: \"Need Python, TypeScript, Rust... not just C#\"",
    "AI: \"TreeSitter has bindings for 30+ languages\"",
    "Human: \"Make one agent that handles all of them\"",
    "",
    "Result: One ingestor, 30+ languages"
], code="""private static readonly Dictionary<string, LanguageConfig> 
    SupportedLanguages = new()
{
    ["py"] = new("python", 
        ["function_definition", "class_definition"]),
    ["ts"] = new("typescript", 
        ["function_declaration", "interface_declaration"]),
    ["rs"] = new("rust", 
        ["function_item", "struct_item"]),
    ["go"] = new("go", 
        ["function_declaration"]),
    // ... 26 more
};""")

# The Pivot - Table
add_table_slide("The Graveyard of Ambition", 
    ["#", "Codename", "Projects", "Lines", "What Happened"],
    [
        ["1", "bird-constellation", "6", "~12k", "Document indexing. Windows-only."],
        ["2", "birdlet", "8", "~18k", "RAG platform. 779-line Program.cs!"],
        ["3", "hve-hack", "17", "~38k", "Full orchestration. Too complex."],
        ["4", "Aura", "4", "~8k", "Local-first foundation. ✓"],
    ])

# Philosophy shift
add_content_slide("The November 25, 2025 Decision", [
    "\"We've built this three times.\"",
    "\"Each time bigger, more complex.\"",
    "",
    "\"What if we kept the patterns and deleted the code?\"",
    "",
    "Philosophy shift:",
    "  FROM: \"How do we orchestrate everything?\"",
    "  TO: \"How do we make each piece excellent?\""
])

# What I Learned
add_content_slide("What I Learned", [
    "AI is a power tool, not autopilot — You still steer",
    "Specs are the leverage point — Control input, not output",
    "Extensibility emerges from dialogue — Talk to your AI",
    "Local-first is about data, not models — Index locally, generate anywhere"
])

# What's Next
add_content_slide("What's Next", [
    "Research module — Document exploration, paper analysis",
    "Personal module — Notes, calendar integration",
    "Self-improvement — Agents that improve agents",
    "",
    "All building on the same Foundation"
])

# Closing
add_title_slide("Questions?", "github.com/johnazariah/aura")

# Save
output_path = r"c:\work\aura\.project\presentations\aura-developer-talk-2025-01.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
